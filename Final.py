import fitz
import os
import asyncio
import aiohttp
from pymongo import MongoClient
from sentence_transformers import SentenceTransformer
from langchain_community.vectorstores import FAISS
from langchain_core.vectorstores import VectorStoreRetriever
from langchain.chains import RetrievalQA
from langchain_community.chat_models.sambanova import ChatSambaNovaCloud
from langchain_core.prompts import PromptTemplate
from langchain_core.embeddings import Embeddings
import faiss
import csv
from docx import Document 
from langchain.memory import ConversationBufferMemory
from langchain.chains.summarize import load_summarize_chain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.document_transformers.embeddings_redundant_filter import EmbeddingsClusteringFilter
import pandas as pd
from pptx import Presentation
import warnings
warnings.filterwarnings('ignore', category=UserWarning)

os.environ['KMP_DUPLICATE_LIB_OK']='TRUE'
API_KEY = os.getenv("SAMBA_API_KEY", "540f8914-997e-46c6-829a-ff76f5d4d265")
API_URL = os.getenv("SAMBA_API_URL", "https://api.sambanova.ai/v1/chat/completions")

client = MongoClient("mongodb+srv://2217028:NquBh8oPPopA0Zuu@sumrag.ux9hs.mongodb.net/?retryWrites=true&w=majority&appName=SUMRAG")
db = client['DT3_EKR_BASE']
collection = db['vectors']

global_faiss_index = None
class SentenceTransformerEmbeddings(Embeddings):
    """Wrapper around SentenceTransformer for compatibility."""
    def __init__(self, model):
        self.model = model

    def embed_documents(self, texts):
        return self.model.encode(texts, show_progress_bar=True, convert_to_tensor=True)

    def embed_query(self, text):
        return self.model.encode(text, convert_to_tensor=True)

embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extracts text from a PDF using PyMuPDF."""
    text = ""
    with fitz.open(pdf_path) as pdf:
        for page_num in range(pdf.page_count):
            page = pdf.load_page(page_num)
            text += page.get_text()
    return text

def extract_text_from_csv(csv_path: str) -> str:
    """Extracts and formats text from a CSV file with a fallback encoding."""
    text = ""
    try:
        with open(csv_path, mode='r', encoding='utf-8') as csv_file:
            reader = csv.reader(csv_file)
            for row in reader:
                text += " ".join(row) + "\n"
    except UnicodeDecodeError:
        with open(csv_path, mode='r', encoding='ISO-8859-1') as csv_file:
            reader = csv.reader(csv_file)
            for row in reader:
                text += " ".join(row) + "\n"
    return text

def extract_text_from_docx(docx_path: str) -> str:
    """Extracts text from a DOCX file."""
    text = ""
    doc = Document(docx_path)
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_excel(excel_path: str) -> str:
    try:
        excel_data = pd.read_excel(excel_path, sheet_name=None, engine='openpyxl')
        sheet_summaries = {}
        Content = ""
        for sheet_name, df in excel_data.items():
            if df.empty:
                sheet_summaries[sheet_name] = "This sheet is empty."
            else:
                
                sheet_content = df.to_string(index=False)
                Content += sheet_content+'\n'
                sheet_summaries[sheet_name] = "Not Empty"
                
        return Content
    except Exception as e:
        raise Exception(f"Failed to extract data from Excel: {e}")

def extract_text_from_ppt(ppt_path):
    text = ""
    try:
        presentation = Presentation(ppt_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting PPT: {e}")
    return text
def extract_text_from_txt(file_path: str) -> str:
    """Extracts text from a plain text file."""
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()
    return text

def chunk_text(text: str, max_length: int = 16000) -> list:
    """Splits the text into larger chunks for faster processing."""
    words = text.split()
    chunks, current_chunk = [], []
    
    for word in words:
        if sum(len(w) for w in current_chunk) + len(word) < max_length:
            current_chunk.append(word)
        else:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]
    
    if current_chunk:
        chunks.append(" ".join(current_chunk))
    
    return chunks


def chunk_text_csv(text: str, max_length: int = 16000) -> list:
    """Splits the text into smaller chunks."""
    words = text.split()
    chunks, current_chunk = [], []
    for word in words:
        if sum(len(w) for w in current_chunk) + len(word) < max_length:
            current_chunk.append(word)
        else:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks

def filter_documents(texts):
    embedding_model1 = SentenceTransformerEmbeddings(model=SentenceTransformer('all-MiniLM-L6-v2',trust_remote_code=True))
    embeddings_filter = EmbeddingsClusteringFilter(
        embeddings=embedding_model1,
        num_clusters=max(1,min(len(texts) // 2, 8)),  
        num_closest=max(1,min(len(texts) // 4, 3)),    
        threshold=0.85                          
    )
    filtered_texts = embeddings_filter.transform_documents(texts)
    return filtered_texts
def extract_text(file):
    file_extension = os.path.splitext(file)[1].lower()
    
    if file_extension == '.pdf':
        loader = PyPDFLoader(file)
        documents = loader.load()
    elif file_extension == '.csv':
        text = extract_text_from_csv(file)
        splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=800)
        documents = splitter.create_documents([text])
    elif file_extension in ['.doc', '.docx']:
        text = extract_text_from_docx(file)
        splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=800)
        documents = splitter.create_documents([text])
    elif file_extension in ['.xlsx', '.xls']:
        text = extract_text_from_excel(file)
        splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=800)
        documents = splitter.create_documents([text])
    elif file_extension in ['.ppt', '.pptx']:
        text = extract_text_from_ppt(file)
        splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=800)
        documents = splitter.create_documents([text])
    elif file_extension == '.txt':
        text = extract_text_from_txt(file)
        splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=800)
        documents = splitter.create_documents([text])
    else:
        raise ValueError("Unsupported file format. Please use PDF, CSV, DOCX, Excel, PPT, or TXT files.")
    
    return documents


# async def summarize_large_text(text: str) -> str:
#     """Optimized async summarization with concurrent processing."""
#     chunks = chunk_text(text)
#     summaries = []
    
#     async with aiohttp.ClientSession() as session:
#         # Process chunks concurrently in batches
#         batch_size = 5  # Process 5 chunks simultaneously
#         for i in range(0, len(chunks), batch_size):
#             batch = chunks[i:i + batch_size]
#             tasks = [summarize_chunk(session, chunk) for chunk in batch]
#             batch_results = await asyncio.gather(*tasks, return_exceptions=True)
            
#             for result in batch_results:
#                 if isinstance(result, str):
#                     summaries.append(result)
    
#     # Combine and create final summary
#     combined_summary = "\n".join(summaries)
    
#     # Get final condensed summary if needed
#     if len(summaries) > 1:
#         final_summary = await summarize_chunk(session, combined_summary)
#         return final_summary
    
#     return combined_summary
# import time

# async def summarize_chunk(session, chunk: str) -> str:
#     # Add delay between API calls
#     time.sleep(1)  # 1 second delay
#     headers = {
#         "Authorization": f"Bearer {API_KEY}",
#         "Content-Type": "application/json"
#     }
#     prompt = f"""Provide a comprehensive summary of the following content in approximately 300-400 words:

#     {chunk}

#     Focus on key points and maintain coherent flow while preserving essential information."""

#     payload = {
#         "model": "llama3-70b",
#         "messages": [{"role": "user", "content": prompt}],
#         "max_tokens": 1024,  
#         "temperature": 0.5,
#         "top_k": 1,
#         "top_p": 0.01
#     }

#     async with session.post(API_URL, json=payload, headers=headers, timeout=30) as response:
#         if response.status == 200:
#             result = await response.json()
#             return result["choices"][0]["message"]["content"]
#         else:
#             raise Exception(f"Error: {response.status} - {await response.text()}")
# async def summarize_large_text(text: str) -> str:
#     """Asynchronously summarizes large text by processing chunks."""
#     chunks = chunk_text(text)
#     summaries = []

#     async with aiohttp.ClientSession() as session:
#         tasks = [summarize_chunk(session, chunk) for chunk in chunks]

#         for task in asyncio.as_completed(tasks):
#             try:
#                 summary = await task
#                 summaries.append(summary)
#             except Exception as e:
#                 pass

#     return "\n".join(summaries)
async def summarize_document(file):
    texts = extract_text(file)
    filtered_docs = filter_documents(texts)
    
    prompt_template = """
    Please provide a comprehensive summary of the following document. Focus on:
    1. A detailed summary of the document
    2. Main themes and key points
    3. Important findings or conclusions
    4. Significant data or statistics if present
    5. Key recommendations or actions

    Document content:
    {text}

    Summary:
    """
    llm1 = ChatSambaNovaCloud(
        model="llama3-70b",
        temperature=0.6,
        max_tokens = 4000
    )
    chain = load_summarize_chain(
        llm1,
        chain_type="stuff",
        prompt=PromptTemplate(template=prompt_template, input_variables=["text"])
    )
    
    summary = chain.invoke(filtered_docs)
    return summary['output_text']

# async def summarize_large_text(text: str) -> str:
#     """Generates a comprehensive overall summary for any input file type"""
#     chunks = chunk_text(text)
#     summaries = []
    
#     async with aiohttp.ClientSession() as session:
#         for chunk in chunks:
#             summary = await summarize_chunk(session, chunk)
#             summaries.append(summary)
#         combined_summary = " ".join(summaries)
#         final_prompt = f"""Create a comprehensive overall summary of this entire document, highlighting the main themes, key points, and important conclusions:

#         {combined_summary}

#         Provide a well-structured summary that captures the essence of the entire document."""

#         payload = {
#             "model": "llama3-70b",
#             "messages": [{"role": "user", "content": final_prompt}],
#             "max_tokens": 1500,
#             "temperature": 0.7,
#             "top_k": 1,
#             "top_p": 0.01
#         }

#         async with session.post(API_URL, json=payload, headers={
#             "Authorization": f"Bearer {API_KEY}",
#             "Content-Type": "application/json"
#         }) as response:
#             if response.status == 200:
#                 result = await response.json()
#                 return result["choices"][0]["message"]["content"]
#             else:
#                 raise Exception(f"Error: {response.status} - {await response.text()}")

def update_global_faiss_index(content_chunks):
    """Update the global FAISS index with new content chunks."""
    global global_faiss_index
    embeddings = SentenceTransformerEmbeddings(embedding_model)
    current_faiss_index = FAISS.from_texts(
        texts=content_chunks,
        embedding=embeddings,
        metadatas=[{"source": "pdf", "index": i} for i in range(len(content_chunks))]
    )

    if global_faiss_index is None:
        global_faiss_index = current_faiss_index
    else:
        global_faiss_index.merge_from(current_faiss_index)
# def get_custom_prompt_for_pdf_and_docx():
    # custom_prompt_template = """<s>[INST] <<SYS>>
    # You are an intelligent assistant designed to extract and provide comprehensive, detailed answers based on the content of a document.
    # Use the following context, extracted from the file, to answer the user's questions. When responding, locate and include all relevant information related to the question from the document.
    # You must NEVER:
    #      - Generate code examples
    #      - Provide general knowledge
    #      - Answer questions outside the document scope
    #      - Make assumptions beyond the text
    # Context:
    # {context}

    # <</SYS>>

    # Please answer the following question:
    # {question}

    # Ensure your response includes all available and relevant content from the document to fully address the user's query. If the exact answer is not found within the data, respond with:
    # 'This information isn't available in the provided data.'
    # [/INST]"""
    # return PromptTemplate(template=custom_prompt_template, input_variables=['context', 'question'])

def get_custom_prompt_for_pdf_and_docx():
      custom_prompt_template = """<s>[INST] <<SYS>>
      You are a focused document expert with these strict rules:
      1. You are ONLY authorized to discuss content from the provided document . You are an intelligent assistant designed to extract and provide comprehensive, detailed answers based on the content of a document
      2. For ANY question not directly addressed in the document, respond with:
         "Let's focus on the document content. Please ask questions about the information in the document."
      3.You must NEVER Generate code examples if asked reponse with:
         "Let's focus on the document content. Please ask questions about the information in the document."
         If the code for the query is available in the provided data , then you provide the code from document only . You should not provide any outside code.
      4.You must NEVER Provide general knowledge if asked reponse with:
         "Let's focus on the document content. Please ask questions about the information in the document."
      5.You must NEVER Answer questions outside the document scope if asked reponse with:
         "Let's focus on the document content. Please ask questions about the information in the document."
      6.You must NEVER Make assumptions beyond the text if asked reponse with:
         "Let's focus on the document content. Please ask questions about the information in the document."
    
      Document Context:
      {context}

      <</SYS>>

      Question: {question}
      Ensure your response includes all available and relevant content from the document to fully address the user's query.
      Provide information ONLY if it exists in the document context above.
      [/INST]"""
      return PromptTemplate(template=custom_prompt_template, input_variables=['context', 'question'])

llm = ChatSambaNovaCloud(
      model="llama3-70b",
      max_tokens=4000,
      temperature=0.5, 
      top_k=1,
      top_p=0.1, 
      request_timeout=45
)

# def get_custom_prompt_csv():
#     custom_prompt_template = """<s>[INST] <<SYS>>
#            You are an advanced information extraction and question-answering assistant designed to provide comprehensive, detailed responses 
#            based strictly on the context from the uploaded CSV file. Your goal is to:

#            1. Thoroughly analyze the entire context
#            2. Extract ALL relevant information related to the user's query
#       3. Provide a comprehensive, multi-faceted response
#     4. If information is partially available, include all related details
#     5. Organize the response in a clear, structured manner
#     6. Be as exhaustive as possible within the context of the provided data

#     Important Guidelines:
#     - If the query can be answered completely or partially from the context, provide a detailed response
#     - Include multiple perspectives or aspects related to the query
#     - If some information is missing, clearly state which parts are covered
#     - Avoid adding any external or hypothetical information
#     - If no information is found, explicitly explain that no relevant information exists in the context

#     Context:
#     {context}

#     <</SYS>>

#     User Query: {question}

#     Detailed Response Requirements:
#     - Provide a comprehensive answer
#     - Break down the response into clear sections if multiple aspects are relevant
#     - Cite specific details from the context
#     - If the information is insufficient, explain exactly what is missing

#     [/INST]"""
#     return PromptTemplate(template=custom_prompt_template, input_variables=['context', 'question'])

def get_custom_prompt_csv():
    custom_prompt_template = """<s>[INST] <<SYS>>
    You are an advanced information extraction and question-answering assistant designed to provide comprehensive, detailed responses 
    based strictly on the context from the uploaded CSV file. Your goal is to:

    1. Thoroughly analyze the entire context
    2. Extract ALL relevant information related to the user's query
    3. Provide a comprehensive, multi-faceted response
    4. If information is partially available, include all related details
    5. Organize the response in a clear, structured manner
    6. Be as exhaustive as possible within the context of the provided data

    Important Guidelines:
    - If the query can be answered completely or partially from the context, provide a detailed response
    - Include multiple perspectives or aspects related to the query
    - If some information is missing, clearly state which parts are covered
    - Avoid adding any external or hypothetical information
    - If no information is found, explicitly explain that no relevant information exists in the context

    Context:
    {context}

    <</SYS>>

    User Query: {question}

    Detailed Response Requirements:
    - Provide a comprehensive answer
    - Break down the response into clear sections if multiple aspects are relevant
    - Cite specific details from the context
    - If the information is insufficient, explain exactly what is missing

    [/INST]"""
    return PromptTemplate(template=custom_prompt_template, input_variables=['context', 'question'])
def get_custom_prompt_excel():
    custom_prompt_template = """<s>[INST] <<SYS>>
    You are an intelligent assistant designed to extract and provide detailed, accurate answers based on the contents of an document.
    Use the following context, extracted from the file, to answer the user's questions. If the question is related to specific data in the document, locate and include all relevant information.

    Excel Context:
    {context}

    <</SYS>>

    Please answer the following question:
    {question}

    Ensure your response is based on precise data from the Excel file and covers all details available in the context. If you cannot find an answer within the data, respond with:
    'This information isn't available in the provided data.'
    [/INST]"""
    return PromptTemplate(template=custom_prompt_template, input_variables=['context', 'question'])

def get_custom_prompt_pptx():
    custom_prompt_template = """<s>[INST] <<SYS>>
    You are an intelligent assistant designed to extract and provide comprehensive, detailed answers based on the content of a document.
    Use the following context, extracted from the file, to answer the user's questions. When responding, locate and include all relevant information related to the question from the document.

    Context:
    {context}

    <</SYS>>

    Please answer the following question:
    {question}

    Ensure your response includes all available and relevant content from the document to fully address the user's query. If the exact answer is not found within the data, respond with:
    'This information isn't available in the provided data.'
    [/INST]"""
    return PromptTemplate(template=custom_prompt_template, input_variables=['context', 'question'])

def save_summary_and_interaction_to_mongodb(pdf_summary, interactions):
    """Save the PDF summary and interactions to MongoDB."""
    document = {
        "pdf_summary": pdf_summary,
        "interactions": interactions
    }
    collection.insert_one(document)
    print("Summary and interactions saved to MongoDB.")

async def main():
    global global_faiss_index  

    while True:
        file_path = input("Enter the file path (or type 'exit' to quit): ").strip()
        if file_path.lower() == 'exit':
            print("Exiting the assistant. Goodbye!")
            break

        file_extension = os.path.splitext(file_path)[1].lower()
        
        try:
            summary = await summarize_document(file_path)
            print("\n" + "="*50)
            print("="*50)
            print(summary)
            print("="*50 + "\n")
            
            if file_extension == '.pdf':
                text = extract_text_from_pdf(file_path)
                custom_prompt = get_custom_prompt_for_pdf_and_docx()
                content_chunks = chunk_text(text)
                print("PDF Text Extracted Successfully.")
            elif file_extension == '.docx':
                text = extract_text_from_docx(file_path)
                custom_prompt = get_custom_prompt_for_pdf_and_docx()
                content_chunks = chunk_text(text)
                print("DOCX Text Extracted Successfully.")
            elif file_extension == '.csv':
                text = extract_text_from_csv(file_path)
                custom_prompt = get_custom_prompt_csv()
                content_chunks = chunk_text_csv(text)
                print("CSV Text Extracted Successfully.")
            elif file_extension in ['.xlsx', '.xls']:
                text = extract_text_from_excel(file_path)
                custom_prompt = get_custom_prompt_excel()
                content_chunks = chunk_text(text)
                print("Excel Text Extracted Successfully.")
            elif file_extension in ['.ppt', '.pptx']:
                text = extract_text_from_ppt(file_path)
                custom_prompt = get_custom_prompt_pptx()
                content_chunks = chunk_text(text)
                print("PPT Text Extracted Successfully.")
            elif file_extension == '.txt':
                text = extract_text_from_txt(file_path)
                custom_prompt = get_custom_prompt_for_pdf_and_docx()
                content_chunks = chunk_text(text)
                print("TXT Text Extracted Sucessfully.")

            else:
                print("Unsupported file format. Please use PDF, DOCX, or CSV files.")
                continue 

            if not content_chunks:
                print("No content chunks generated. Please check the file content.")
                continue
            update_global_faiss_index(content_chunks)
            print("Global FAISS index updated successfully.")


            retriever = VectorStoreRetriever(vectorstore=global_faiss_index, search_kwargs={"k": 5,"fetch_k":10})
            llm = ChatSambaNovaCloud(
                model="llama3-70b",
                max_tokens=4000,
                temperature=0.5,
                top_k=1,
                top_p=0.8,
                request_timeout=30
            )

            memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)
            qa_chain = RetrievalQA.from_chain_type(
                llm=llm,
                retriever=retriever,
                chain_type_kwargs={"prompt": custom_prompt, "document_variable_name": "context"},
                memory=memory
            )

            interactions = []
            while True:
                user_query = input("Ask a question (or type 'next' for another file, 'exit' to quit): ").strip()
                if user_query.lower() == 'exit':
                    print("Exiting the assistant. Goodbye!")
                    return
                elif user_query.lower() == 'next':
                    print("Processing the next file.")
                    break

                try:
                    result = qa_chain.invoke({"query": user_query})
                    response = result['result']
                    print(f"Response: {response}")
                    interactions.append({"query": user_query, "response": response})
                except Exception as e:
                    print(f"Error during query: {e}")

        except Exception as e:
            print(f"Failed to process file: {e}")
            continue
        # save_summary_and_interaction_to_mongodb(summary, interactions)
if __name__ == "__main__":
    asyncio.run(main())




